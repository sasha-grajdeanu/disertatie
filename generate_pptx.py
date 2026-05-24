#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Colors ───
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE   = RGBColor(0x16, 0x21, 0x3E)
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2    = RGBColor(0x7B, 0x68, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
CODE_BG    = RGBColor(0x0D, 0x11, 0x17)
ORANGE     = RGBColor(0xFF, 0xA5, 0x00)
GREEN      = RGBColor(0x00, 0xE6, 0x76)
RED        = RGBColor(0xFF, 0x45, 0x45)
TABLE_HEAD = RGBColor(0x1E, 0x3A, 0x5F)
TABLE_ROW1 = RGBColor(0x14, 0x1E, 0x30)
TABLE_ROW2 = RGBColor(0x1A, 0x26, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    return p

def add_bullet(tf, text, size=16, color=LIGHT_GRAY, indent=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_before = Pt(4)
    p.level = indent
    return p

def add_code_box(slide, left, top, width, height, lines):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = GREEN
        p.font.name = "Consolas"
        p.space_before = Pt(2)
    return tf

def add_accent_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(2), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def make_table(slide, left, top, width, rows_data, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.4 * rows))
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT
                else:
                    p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEAD
            elif r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW1
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW2
    return table

# ═══════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text(slide, 1, 1.8, 11, 1.2, "VENUS", size=60, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.0, 11, 0.8, "Automated Test Generation via Rewriting Logic", size=28, color=WHITE, align=PP_ALIGN.CENTER)

tf = add_text(slide, 1, 4.5, 11, 0.5, "Sasha Grajdeanu", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_para(tf, "Dissertation Project", size=16, color=LIGHT_GRAY)

# accent bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2.85), Inches(3.333), Inches(0.05))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# ═══════════════════════════════════════════
# SLIDE 2 — Motivation & Approach
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)

add_text(slide, 0.8, 0.4, 10, 0.6, "Motivation & Approach", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

# Left column — Problem
tf = add_text(slide, 0.8, 1.3, 5.5, 0.5, "The Problem", size=22, color=ORANGE, bold=True)
add_bullet(tf, "Writing unit tests is manual and tedious")
add_bullet(tf, "Developers must guess which inputs matter")
add_bullet(tf, "Edge cases and runtime errors are easy to miss")

# Right column — Solution
tf = add_text(slide, 7, 1.3, 5.5, 0.5, "Our Solution", size=22, color=GREEN, bold=True)
add_bullet(tf, "Venus = imperative language + Maude formal engine")
add_bullet(tf, "Execution defined as term rewriting (Rewriting Logic)")
add_bullet(tf, "Same engine executes code OR explores all paths")

# Architecture
tf = add_text(slide, 0.8, 4.2, 11.5, 0.5, "Architecture", size=20, color=ACCENT, bold=True)
add_code_box(slide, 0.8, 4.8, 11.5, 1.2, [
    ".venus file  ──►  Python CLI  ──►  Maude Engine  ──►  Results",
    "                                    ├─ VENUS-TS    (normal execution)",
    "                                    └─ VENUS-TEST  (test discovery)",
])

tf = add_text(slide, 0.8, 6.2, 11.5, 0.5, "python main.py program.venus        →  Execute normally", size=14, color=LIGHT_GRAY)
add_para(tf, "python main.py program.venus -s   →  Discover all test scenarios", size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════
# SLIDE 3 — The Formal Machine
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)

add_text(slide, 0.8, 0.4, 10, 0.6, "The Formal Machine", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

tf = add_text(slide, 0.8, 1.3, 11, 0.5, "State  =  < Control  |  Stack  |  Memory >", size=24, color=WHITE, bold=True, font_name="Consolas")
add_bullet(tf, "Control — instruction stream (what to execute next)")
add_bullet(tf, "Stack — intermediate values during expression evaluation")
add_bullet(tf, "Memory — variable bindings (associative-commutative)")

# AC Memory
tf = add_text(slide, 0.8, 3.2, 5.3, 0.5, "AC Memory — O(1) variable lookup", size=18, color=GREEN, bold=True)
add_code_box(slide, 0.8, 3.7, 5.3, 1.0, [
    "op _,_ : Memory Memory -> Memory",
    "         [assoc comm id: empty] .",
    "eq get(((X -> N), M), X) = N .",
])

# Desugaring
tf = add_text(slide, 6.8, 3.2, 5.7, 0.5, "Desugaring — minimal core", size=18, color=GREEN, bold=True)
add_code_box(slide, 6.8, 3.7, 5.7, 1.0, [
    "eq for(Init,Cond,Step) do Body od =",
    "   Init ; while Cond do Body;Step od .",
    "eq if C then B fi = if C then B else skip fi .",
])

tf = add_text(slide, 0.8, 5.2, 11, 0.5, "The engine only implements while and binary if — everything else is syntactic sugar.", size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════
# SLIDE 4 — The Testing Engine
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)

add_text(slide, 0.8, 0.4, 10, 0.6, "The Testing Engine — cAutoTest", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

# Step 1
tf = add_text(slide, 0.8, 1.3, 5.5, 0.4, "① Extract constants from AST", size=18, color=ORANGE, bold=True)
add_bullet(tf, "if ('score gte 50)  →  extracts 50")

# Step 2
tf = add_text(slide, 0.8, 2.3, 5.5, 0.4, "② Generate boundary values", size=18, color=ORANGE, bold=True)
add_bullet(tf, "Each N  →  {N-1, N, N+1}")
add_bullet(tf, "Always includes 0 and 1 as base cases")
add_bullet(tf, "50  →  {0, 1, 49, 50, 51}")

# Step 3
tf = add_text(slide, 7, 1.3, 5.5, 0.4, "③ Remove duplicates (AC IntSet)", size=18, color=ORANGE, bold=True)
add_code_box(slide, 7, 1.9, 5.5, 0.7, [
    "eq N ;; N = N .",
    "--- Algebraic dedup, no filtering code",
])

# Step 4
tf = add_text(slide, 7, 3.0, 5.5, 0.4, "④ Nondeterministic execution", size=18, color=ORANGE, bold=True)
add_code_box(slide, 7, 3.6, 5.5, 0.7, [
    "rl [pickValue] :",
    "   testValues(N ;; SSet) => N .",
])

tf = add_text(slide, 0.8, 4.8, 11.5, 0.5, "Maude's search explores every combination simultaneously — all paths tested in one command.", 
              size=17, color=WHITE, bold=True)

# The full rule
add_code_box(slide, 0.8, 5.5, 11.5, 1.5, [
    "rl [DoAutoTest] :",
    "  < cAutoTest # CTRL | STK | (F => [L, S]) , MRY >",
    "  =>",
    "  < (call F (genTestArgs(L,",
    "      0 ;; 1 ;; toIntSet(genBoundaries(extractConstants(S))))))",
    "    # CTRL | STK | (F => [L, S]) , tested(F) , MRY > .",
])


# ═══════════════════════════════════════════
# SLIDE 5 — Demo: Branching
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)

add_text(slide, 0.8, 0.4, 10, 0.6, "Demo: Branching — checkStatus", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

add_code_box(slide, 0.8, 1.3, 11.5, 1.6, [
    "def 'checkStatus ('score) {",
    "    'passed := 0 ;",
    "    if ('score gte 50) then",
    "        'passed := 1",
    "    fi ;",
    "    'checked := 1",
    "}",
])

add_text(slide, 0.8, 3.2, 11, 0.4, "Result: 2 execution paths discovered", size=22, color=GREEN, bold=True)

make_table(slide, 0.8, 3.8, 11.5, [
    ["Scenario", "Inputs", "Output"],
    ["1", "'score = 0,  'score = 1,  'score = 49", "'passed = 0,  'checked = 1"],
    ["2", "'score = 50,  'score = 51", "'passed = 1,  'checked = 1"],
], col_widths=[1.5, 5.5, 4.5])

tf = add_text(slide, 0.8, 5.3, 11, 0.5, "✓  Boundary at 50 found automatically — zero manual tests written", size=18, color=GREEN, bold=True)


# ═══════════════════════════════════════════
# SLIDE 6 — Demo: Loops
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)

add_text(slide, 0.8, 0.4, 10, 0.6, "Demo: Loops — classify", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

add_code_box(slide, 0.8, 1.3, 11.5, 1.6, [
    "def 'classify ('n) {",
    "    'sum := 0 ; 'i := 1 ;",
    "    while ('i lte 'n) do 'sum := 'sum + 'i ; 'i := 'i + 1 od ;",
    "    if ('sum gt 10) then 'category := 2",
    "    else if ('sum gt 0) then 'category := 1",
    "         else 'category := 0 fi fi",
    "}",
])

tf = add_text(slide, 0.8, 3.1, 11, 0.4, "Constants: 0, 1, 2, 10  →  Boundaries: -1, 0, 1, 2, 3, 9, 10, 11", size=16, color=LIGHT_GRAY)
add_text(slide, 0.8, 3.5, 11, 0.4, "Result: 5 execution paths discovered", size=22, color=GREEN, bold=True)

make_table(slide, 0.8, 4.1, 11.5, [
    ["Scenario", "Input 'n", "Loop iterations", "'sum", "'category"],
    ["1", "0, -1", "0", "0", "0"],
    ["2", "1", "1", "1", "1"],
    ["3", "9", "9", "45", "2"],
    ["4", "10", "10", "55", "2"],
    ["5", "11", "11", "66", "2"],
], col_widths=[1.5, 2.2, 2.5, 2.3, 2.5])

tf = add_text(slide, 0.8, 6.7, 11, 0.5, "✓  Loop + nested conditionals explored automatically, up to 11 iterations", size=18, color=GREEN, bold=True)


# ═══════════════════════════════════════════
# SLIDE 7 — Demo: Runtime Errors
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)

add_text(slide, 0.8, 0.4, 10, 0.6, "Demo: Runtime Error Detection", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

add_code_box(slide, 0.8, 1.3, 5, 0.8, [
    "def 'divide ('x , 'y) {",
    "    'result := 'x / 'y",
    "}",
])

# How it works
tf = add_text(slide, 7, 1.3, 5.5, 0.4, "How it works:", size=20, color=WHITE, bold=True)
add_bullet(tf, "Division rule has a guard: N2 ≠ 0")
add_bullet(tf, "When 'y = 0, the rule cannot fire")
add_bullet(tf, "State gets stuck — no crash")
add_bullet(tf, "Maude safely isolates the error branch")

# Error output
add_code_box(slide, 0.8, 3.0, 11.5, 2.2, [
    "⚠  STUCK STATES (Runtime Errors): 2",
    "─────────────────────────────────────",
    "Error 1:",
    "   Input: 'x = 0, 'y = 0",
    "   Stuck at: cDiv # cAssign",
    "   Cause: Division by zero",
    "",
    "Error 2:",
    "   Input: 'x = 1, 'y = 0",
    "   Stuck at: cDiv # cAssign",
    "   Cause: Division by zero",
])

# guard rule
tf = add_text(slide, 0.8, 5.5, 11.5, 0.4, "The guard in the Maude rule:", size=16, color=LIGHT_GRAY)
add_code_box(slide, 0.8, 5.9, 11.5, 0.7, [
    "crl [DivEval] : < cDiv # CTRL | N2 $ N1 $ STK | MRY >",
    "  => < CTRL | (N1 quo N2) $ STK | MRY >  if (N2 =/= 0) .",
])


# ═══════════════════════════════════════════
# SLIDE 8 — Future Directions
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)

add_text(slide, 0.8, 0.4, 10, 0.6, "Future Directions", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

tf = add_text(slide, 0.8, 1.5, 11.5, 0.5, "Optimize the code", size=24, color=ORANGE, bold=True)
add_bullet(tf, "Reduce redundant rewrite paths in the Maude search space", size=18)
add_bullet(tf, "Improve performance for programs with larger state spaces", size=18)

tf = add_text(slide, 0.8, 3.0, 11.5, 0.5, "Handle more edge cases", size=24, color=ORANGE, bold=True)
add_bullet(tf, "Test with more complex program structures (deeply nested conditionals, multi-parameter functions)", size=18)
add_bullet(tf, "Identify and fix cases not yet covered by the testing engine", size=18)

tf = add_text(slide, 0.8, 4.5, 11.5, 0.5, "Extend the language", size=24, color=ORANGE, bold=True)
add_bullet(tf, "Add new Venus constructs and verify the testing engine adapts correctly", size=18)
add_bullet(tf, "Expand the set of example programs to validate robustness", size=18)


# ═══════════════════════════════════════════
# SLIDE 9 — Key Takeaways
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, 0.8, 0.4, 10, 0.6, "Key Takeaways", size=32, color=ACCENT, bold=True)
add_accent_line(slide, 1.0)

takeaways = [
    ("1.", "Rewriting Logic is a practical execution engine, not just theory"),
    ("2.", "Nondeterminism gives exhaustive path exploration for free"),
    ("3.", "AC matching eliminates boilerplate (memory lookup, deduplication)"),
    ("4.", "The testing framework requires zero manual test definitions"),
    ("5.", "Works on branching, loops, and error detection — all automatically"),
]

y = 1.5
for num, text in takeaways:
    tf = add_text(slide, 1.2, y, 1, 0.5, num, size=22, color=ACCENT, bold=True)
    add_text(slide, 1.8, y, 10, 0.5, text, size=20, color=WHITE)
    y += 0.65

# Quote
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(10.333), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x10, 0x18, 0x28)
shape.line.color.rgb = ACCENT2
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The same formal semantics that define the language also test it — automatically and exhaustively."
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.italic = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

add_text(slide, 1, 6.3, 11.333, 0.6, "Thank you!  Questions?", size=30, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ─── Save ───
prs.save("Venus_Presentation.pptx")
print("✓ Saved Venus_Presentation.pptx")
